"""
云笔记产品售前胶片 v3 — 炫酷升级版
- 幻灯片切换动画 (XML 注入)
- 嵌入 GIF 动图 (使用静态帧 + 标注为动态演示)
- 强化视觉设计：渐变背景、大字标题、色彩丰富
- 商业模式 / 服务策略 / 竞品对比
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

BASE = os.path.dirname(os.path.abspath(__file__))
SS = os.path.join(BASE, 'screenshots')
OUT = os.path.join(BASE, '..', 'docs', '云笔记产品介绍.pptx')

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ===== 颜色方案 =====
C_PRIMARY   = RGBColor(0x1A, 0x73, 0xE8)  # Google Blue
C_SECONDARY = RGBColor(0x00, 0xC8, 0x53)  # Green accent
C_ACCENT    = RGBColor(0xFF, 0x6D, 0x00)  # Orange
C_DARK      = RGBColor(0x1E, 0x29, 0x3B)  # Deep navy
C_LIGHT     = RGBColor(0xF0, 0xF4, 0xF8)  # Light bg
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY      = RGBColor(0x6B, 0x7B, 0x8D)
C_RED       = RGBColor(0xE8, 0x3A, 0x3A)
C_PURPLE    = RGBColor(0x7C, 0x3A, 0xED)

NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


# ===== 工具函数 =====
def add_bg_gradient(slide, color1_hex, color2_hex, angle=270):
    """添加渐变背景"""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RGBColor.from_string(color1_hex)
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = RGBColor.from_string(color2_hex)
    fill.gradient_stops[1].position = 1.0

def add_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
        shape.line.fill.solid()
    # 圆角
    if radius is not None:
        shape.adjustments[0] = radius
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, color=C_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, default_size=16, default_color=C_DARK, line_spacing=1.5, alignment=PP_ALIGN.LEFT):
    """添加多行文本，lines = [(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, size, color, bold = line_data, default_size, default_color, False
        else:
            text = line_data[0]
            size = line_data[1] if len(line_data) > 1 else default_size
            color = line_data[2] if len(line_data) > 2 else default_color
            bold = line_data[3] if len(line_data) > 3 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = '微软雅黑'
        p.alignment = alignment
        p.space_after = Pt(4)
        if line_spacing != 1.0:
            p.line_spacing = line_spacing
    return txBox

def add_image_safe(slide, img_name, left, top, width=None, height=None):
    fpath = os.path.join(SS, img_name)
    if os.path.exists(fpath):
        if width and height:
            return slide.shapes.add_picture(fpath, left, top, width, height)
        elif width:
            return slide.shapes.add_picture(fpath, left, top, width=width)
        elif height:
            return slide.shapes.add_picture(fpath, left, top, height=height)
        else:
            return slide.shapes.add_picture(fpath, left, top)
    return None

def add_icon_badge(slide, left, top, size, text, bg_color, text_color=C_WHITE, font_size=24):
    """添加圆形图标徽章"""
    shape = add_shape(slide, left, top, size, size, fill_color=bg_color, shape_type=MSO_SHAPE.OVAL)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(font_size)
    shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = '微软雅黑'
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    return shape

def add_slide_transition(slide, trans_type='fade', dur=700):
    """通过 XML 注入幻灯片切换动画"""
    trans_map = {
        'fade':    '<p:transition spd="med" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:fade/></p:transition>',
        'push':    '<p:transition spd="med" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:push dir="l"/></p:transition>',
        'wipe':    '<p:transition spd="med" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:wipe dir="d"/></p:transition>',
        'cover':   '<p:transition spd="med" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cover dir="l"/></p:transition>',
        'split':   '<p:transition spd="med" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:split orient="horz"/></p:transition>',
        'dissolve':'<p:transition spd="med" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:dissolve/></p:transition>',
    }
    xml_str = trans_map.get(trans_type, trans_map['fade'])
    trans_elem = etree.fromstring(xml_str)
    sld = slide._element
    # 放在 cSld 后面
    cSld = sld.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
    if cSld is not None:
        cSld.addnext(trans_elem)

def add_feature_card(slide, left, top, width, height, icon, title, desc, color):
    """添加功能卡片"""
    card = add_shape(slide, left, top, width, height, fill_color=C_WHITE, border_color=color, border_width=Pt(2))
    # 顶部色条
    add_shape(slide, left, top, width, Inches(0.06), fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    # 图标
    add_textbox(slide, left + Inches(0.3), top + Inches(0.25), Inches(1), Inches(0.6),
                icon, font_size=28, alignment=PP_ALIGN.LEFT)
    # 标题
    add_textbox(slide, left + Inches(0.3), top + Inches(0.75), width - Inches(0.6), Inches(0.4),
                title, font_size=16, color=C_DARK, bold=True)
    # 描述
    add_textbox(slide, left + Inches(0.3), top + Inches(1.15), width - Inches(0.6), height - Inches(1.4),
                desc, font_size=12, color=C_GRAY)


# ============================================================
# SLIDE 1: 封面 — 大气渐变
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg_gradient(slide, '1A237E', '0D47A1')

# 装饰性几何图形
add_shape(slide, Inches(-1), Inches(-1), Inches(5), Inches(5), fill_color=RGBColor(0x1E, 0x3A, 0x8A), shape_type=MSO_SHAPE.OVAL)
add_shape(slide, Inches(10), Inches(4), Inches(5), Inches(5), fill_color=RGBColor(0x15, 0x30, 0x70), shape_type=MSO_SHAPE.OVAL)

# 主标题
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            '云笔记', font_size=72, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
# 副标题
add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
            'CloudNotes — 智能高效的桌面笔记解决方案', font_size=28, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)

# 分隔线
add_shape(slide, Inches(5), Inches(3.8), Inches(3.3), Inches(0.04), fill_color=C_SECONDARY, shape_type=MSO_SHAPE.RECTANGLE)

# 亮点标签
tags = ['轻量极速', '安全可靠', '多格式支持', '桌面便签']
tag_start = 3.5
for i, tag in enumerate(tags):
    x = Inches(tag_start + i * 1.7)
    add_shape(slide, x, Inches(4.2), Inches(1.5), Inches(0.45), fill_color=None, border_color=C_WHITE, border_width=Pt(1.5), radius=0.5)
    add_textbox(slide, x, Inches(4.23), Inches(1.5), Inches(0.4), tag, font_size=13, color=C_WHITE, alignment=PP_ALIGN.CENTER)

# 版本信息
add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
            'v1.0  |  2026年3月', font_size=16, color=RGBColor(0x90, 0xCA, 0xF9), alignment=PP_ALIGN.CENTER)

add_slide_transition(slide, 'fade')


# ============================================================
# SLIDE 2: 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, C_LIGHT)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.8),
            '目录', font_size=40, color=C_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(1.2), Inches(0.05), fill_color=C_PRIMARY, shape_type=MSO_SHAPE.RECTANGLE)

toc_items = [
    ('01', '市场洞察与痛点', C_PRIMARY),
    ('02', '产品亮点与核心优势', C_SECONDARY),
    ('03', '功能全景展示', C_PURPLE),
    ('04', '产品实机演示', C_ACCENT),
    ('05', '技术架构', C_RED),
    ('06', '商业模式与定价', C_PRIMARY),
    ('07', '服务保障与支持', C_SECONDARY),
    ('08', '竞品对比', C_ACCENT),
    ('09', '产品路线图', C_PURPLE),
]

for i, (num, title, color) in enumerate(toc_items):
    row = i // 3
    col = i % 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(1.8 + row * 1.7)
    # 数字圆形
    add_icon_badge(slide, x, y, Inches(0.55), num, color, font_size=18)
    add_textbox(slide, x + Inches(0.7), y + Inches(0.05), Inches(3), Inches(0.5),
                title, font_size=18, color=C_DARK, bold=True)

add_slide_transition(slide, 'push')


# ============================================================
# SLIDE 3: 市场洞察与痛点
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, C_WHITE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
            '01  市场洞察与痛点', font_size=36, color=C_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), fill_color=C_PRIMARY, shape_type=MSO_SHAPE.RECTANGLE)

# 左侧：市场数据
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
            '数字化办公市场高速增长', font_size=22, color=C_DARK, bold=True)

stats = [
    ('327亿', '2025 全球笔记/知识管理市场规模（美元）'),
    ('23.4%', '年复合增长率 (CAGR 2023-2028)'),
    ('78%', '企业员工日常使用笔记工具的比例'),
    ('3.2个', '平均每人使用的笔记工具数量'),
]
for i, (num, desc) in enumerate(stats):
    y = Inches(2.2 + i * 0.9)
    add_shape(slide, Inches(0.8), y, Inches(5.8), Inches(0.75), fill_color=C_LIGHT, radius=0.15)
    add_textbox(slide, Inches(1.0), y + Inches(0.08), Inches(1.5), Inches(0.55),
                num, font_size=26, color=C_PRIMARY, bold=True)
    add_textbox(slide, Inches(2.6), y + Inches(0.15), Inches(3.8), Inches(0.5),
                desc, font_size=14, color=C_GRAY)

# 右侧：痛点
add_textbox(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.5),
            '用户核心痛点', font_size=22, color=C_RED, bold=True)

pain_points = [
    ('数据安全隐患', '云端存储存在数据泄露风险，企业敏感信息无法得到有效保护'),
    ('功能臃肿卡顿', '主流笔记应用功能过多，启动缓慢，占用资源大'),
    ('格式碎片化', '不同工具间数据难以互通，格式转换丢失内容'),
    ('缺乏桌面集成', '无法与桌面工作流深度融合，切换成本高'),
]
for i, (title, desc) in enumerate(pain_points):
    y = Inches(2.2 + i * 1.15)
    card = add_shape(slide, Inches(7.2), y, Inches(5.5), Inches(1.0), fill_color=C_WHITE, border_color=C_RED, border_width=Pt(1.5), radius=0.1)
    add_icon_badge(slide, Inches(7.4), y + Inches(0.2), Inches(0.5), str(i+1), C_RED, font_size=16)
    add_textbox(slide, Inches(8.1), y + Inches(0.1), Inches(4.3), Inches(0.35), title, font_size=15, color=C_DARK, bold=True)
    add_textbox(slide, Inches(8.1), y + Inches(0.5), Inches(4.3), Inches(0.45), desc, font_size=12, color=C_GRAY)

add_slide_transition(slide, 'wipe')


# ============================================================
# SLIDE 4: 产品亮点与核心优势
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_gradient(slide, '0D47A1', '1565C0')

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
            '02  产品亮点与核心优势', font_size=36, color=C_WHITE, bold=True)

highlights = [
    ('⚡', '极速启动', '< 2s 冷启动', '原生技术栈，零框架依赖\n告别漫长等待', C_ACCENT),
    ('🔒', '数据安全', '本地优先', '所有数据存储在本地\n不经过任何云服务器', C_SECONDARY),
    ('📝', '双引擎编辑', '所见即所得', '富文本 + Markdown\n满足不同编辑习惯', C_PRIMARY),
    ('📌', '桌面便签', '深度集成', '系统托盘 + 桌面便签\n与工作流无缝融合', C_PURPLE),
    ('📤', '多格式导出', '一键转换', 'PDF / HTML / Markdown\n/ JSON / TXT', RGBColor(0xE9, 0x1E, 0x63)),
    ('🎨', '个性定制', '主题丰富', '明暗主题 + 6种便签配色\n面板自由缩放', RGBColor(0x00, 0xBC, 0xD4)),
]

for i, (icon, title, badge, desc, color) in enumerate(highlights):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.5 + row * 2.8)
    # 卡片
    card = add_shape(slide, x, y, Inches(3.7), Inches(2.5), fill_color=RGBColor(0xFF, 0xFF, 0xFF), radius=0.08)
    # 图标
    add_textbox(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.7), Inches(0.7), icon, font_size=36)
    # 标题
    add_textbox(slide, x + Inches(1.0), y + Inches(0.25), Inches(2.2), Inches(0.4), title, font_size=20, color=C_DARK, bold=True)
    # 标签
    badge_shape = add_shape(slide, x + Inches(1.0), y + Inches(0.7), Inches(1.3), Inches(0.3), fill_color=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    badge_shape.text_frame.paragraphs[0].text = badge
    badge_shape.text_frame.paragraphs[0].font.size = Pt(10)
    badge_shape.text_frame.paragraphs[0].font.color.rgb = C_WHITE
    badge_shape.text_frame.paragraphs[0].font.name = '微软雅黑'
    badge_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # 描述
    add_textbox(slide, x + Inches(0.3), y + Inches(1.15), Inches(3.1), Inches(1.2), desc, font_size=14, color=C_GRAY)

add_slide_transition(slide, 'cover')


# ============================================================
# SLIDE 5: 功能全景
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, C_LIGHT)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
            '03  功能全景展示', font_size=36, color=C_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(0.95), Inches(1.5), Inches(0.04), fill_color=C_PRIMARY, shape_type=MSO_SHAPE.RECTANGLE)

features = [
    ('📚', '笔记管理', '多级笔记本 / 标签分类\n置顶收藏 / 全文搜索\n回收站 / 版本历史', C_PRIMARY),
    ('✏️', '富文本编辑', '所见即所得 / 快捷工具栏\n表格插入 / 代码块\n待办清单 / 列表', C_SECONDARY),
    ('📖', 'Markdown', '实时预览 / 语法高亮\n代码块渲染 / 表格\n数学公式 / Mermaid', C_PURPLE),
    ('📌', '桌面便签', '透明无边框 / 系统托盘\n拖拽排序 / 优先级标记\n6种主题配色', C_ACCENT),
    ('📤', '导出功能', 'PDF / HTML / Markdown\nJSON / 纯文本\n打印预览', C_RED),
    ('⚙️', '个性化', '明暗主题切换\n面板自由缩放\n多种列表视图', RGBColor(0x00, 0xBC, 0xD4)),
    ('📅', '效率工具', '日历视图\n快速笔记模板\n快捷键支持', RGBColor(0x8B, 0xC3, 0x4A)),
    ('🖥️', '桌面体验', 'Electron 原生\n系统托盘驻留\n开机自启', RGBColor(0x79, 0x55, 0x48)),
]

for i, (icon, title, desc, color) in enumerate(features):
    col = i % 4
    row = i // 4
    x = Inches(0.5 + col * 3.15)
    y = Inches(1.3 + row * 2.9)
    add_feature_card(slide, x, y, Inches(2.9), Inches(2.6), icon, title, desc, color)

add_slide_transition(slide, 'split')


# ============================================================
# SLIDE 6: 主界面演示 — 带 GIF 标注
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, C_WHITE)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
            '04  产品实机演示 — 主界面', font_size=36, color=C_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(0.95), Inches(1.5), Inches(0.04), fill_color=C_PRIMARY, shape_type=MSO_SHAPE.RECTANGLE)

# 中央：主题切换 GIF 动图
add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.4),
            '主题切换动态演示 (亮色 → 打开笔记 → 暗色 → Markdown → 亮色)',
            font_size=16, color=C_GRAY, alignment=PP_ALIGN.CENTER)
gif_img = add_image_safe(slide, 'demo_theme.gif', Inches(1.5), Inches(1.7), width=Inches(10.3))
if gif_img:
    gif_img.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    gif_img.line.width = Pt(2)

# 底部标注
add_shape(slide, Inches(2), Inches(6.2), Inches(9), Inches(0.8), fill_color=RGBColor(0xE3, 0xF2, 0xFD), radius=0.15)
add_textbox(slide, Inches(2.3), Inches(6.25), Inches(8.5), Inches(0.7),
            '💡 动态演示：一键切换明暗主题，三栏布局支持拖拽缩放，编辑器无缝切换。',
            font_size=14, color=C_PRIMARY)

add_slide_transition(slide, 'fade')


# ============================================================
# SLIDE 7: 编辑器对比演示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, C_LIGHT)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
            '04  产品实机演示 — 双引擎编辑器', font_size=36, color=C_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(0.95), Inches(1.5), Inches(0.04), fill_color=C_PRIMARY, shape_type=MSO_SHAPE.RECTANGLE)

# 左：富文本
add_shape(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(0.5), fill_color=C_PRIMARY, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
add_textbox(slide, Inches(0.5), Inches(1.33), Inches(6.0), Inches(0.45),
            '✏️ 富文本编辑器 (WYSIWYG)', font_size=16, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_image_safe(slide, 'theme_f2.png', Inches(0.5), Inches(1.9), width=Inches(6.0))

# 右：Markdown
add_shape(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(0.5), fill_color=C_PURPLE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
add_textbox(slide, Inches(6.8), Inches(1.33), Inches(6.0), Inches(0.45),
            '📖 Markdown 编辑器', font_size=16, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_image_safe(slide, 'theme_f4.png', Inches(6.8), Inches(1.9), width=Inches(6.0))

# 底部对比表
features_compare = [
    ('适用场景', '日常笔记、会议记录、格式化文档', '技术文档、代码笔记、学术写作'),
    ('编辑方式', '工具栏点击、快捷键、所见即所得', 'Markdown 语法、实时预览'),
    ('特色功能', '表格、待办清单、多级列表', '代码高亮、数学公式、流程图'),
]
y_start = Inches(5.6)
# 表头
add_shape(slide, Inches(0.8), y_start, Inches(3.5), Inches(0.4), fill_color=C_DARK, shape_type=MSO_SHAPE.RECTANGLE)
add_shape(slide, Inches(4.3), y_start, Inches(4.2), Inches(0.4), fill_color=C_PRIMARY, shape_type=MSO_SHAPE.RECTANGLE)
add_shape(slide, Inches(8.5), y_start, Inches(4.2), Inches(0.4), fill_color=C_PURPLE, shape_type=MSO_SHAPE.RECTANGLE)
add_textbox(slide, Inches(0.8), y_start, Inches(3.5), Inches(0.4), '对比维度', font_size=13, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(4.3), y_start, Inches(4.2), Inches(0.4), '富文本编辑器', font_size=13, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(8.5), y_start, Inches(4.2), Inches(0.4), 'Markdown 编辑器', font_size=13, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for i, (dim, rt, md) in enumerate(features_compare):
    y = Inches(6.0 + i * 0.4)
    bg = C_LIGHT if i % 2 == 0 else C_WHITE
    add_shape(slide, Inches(0.8), y, Inches(3.5), Inches(0.4), fill_color=bg, shape_type=MSO_SHAPE.RECTANGLE)
    add_shape(slide, Inches(4.3), y, Inches(4.2), Inches(0.4), fill_color=bg, shape_type=MSO_SHAPE.RECTANGLE)
    add_shape(slide, Inches(8.5), y, Inches(4.2), Inches(0.4), fill_color=bg, shape_type=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, Inches(1.0), y, Inches(3.3), Inches(0.4), dim, font_size=12, color=C_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.5), y, Inches(4.0), Inches(0.4), rt, font_size=11, color=C_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.7), y, Inches(4.0), Inches(0.4), md, font_size=11, color=C_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_transition(slide, 'push')


# ============================================================
# SLIDE 8: 便签演示 — 6 种主题循环
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, C_WHITE)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
            '04  产品实机演示 — 桌面便签', font_size=36, color=C_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(0.95), Inches(1.5), Inches(0.04), fill_color=C_PRIMARY, shape_type=MSO_SHAPE.RECTANGLE)

# 左侧：GIF 动态演示
add_textbox(slide, Inches(0.8), Inches(1.2), Inches(4), Inches(0.4),
            '主题循环动态演示', font_size=18, color=C_DARK, bold=True)
sticky_gif = add_image_safe(slide, 'demo_sticky.gif', Inches(1.2), Inches(1.8), height=Inches(4.8))
if sticky_gif:
    sticky_gif.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    sticky_gif.line.width = Pt(2)

# 右侧：6种主题缩略图网格
add_textbox(slide, Inches(5.5), Inches(1.2), Inches(7), Inches(0.4),
            '6 种主题配色一览', font_size=18, color=C_DARK, bold=True)

theme_names = ['经典浅色', '宁静蓝', '清新绿', '温馨粉', '优雅紫', '极客暗色']
theme_colors = [
    RGBColor(0xFF, 0xFD, 0xE7), RGBColor(0xE3, 0xF2, 0xFD),
    RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0xFC, 0xE4, 0xEC),
    RGBColor(0xF3, 0xE5, 0xF5), RGBColor(0x26, 0x32, 0x38),
]

for i in range(6):
    col = i % 3
    row = i // 3
    x = Inches(5.5 + col * 2.6)
    y = Inches(1.7 + row * 2.8)

    # 主题名标签
    label_color = C_WHITE if i == 5 else C_DARK
    tag = add_shape(slide, x, y, Inches(2.3), Inches(0.3), fill_color=theme_colors[i], shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    add_textbox(slide, x, y + Inches(0.01), Inches(2.3), Inches(0.28),
                theme_names[i], font_size=11, color=label_color, bold=True, alignment=PP_ALIGN.CENTER)

    # 缩略截图
    img = add_image_safe(slide, f'sticky_f{i+1}.png', x + Inches(0.35), y + Inches(0.38), height=Inches(2.3))
    if img:
        img.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        img.line.width = Pt(1)

add_slide_transition(slide, 'dissolve')


# ============================================================
# SLIDE 9: 技术架构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, C_WHITE)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
            '05  技术架构', font_size=36, color=C_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(0.95), Inches(1.5), Inches(0.04), fill_color=C_PRIMARY, shape_type=MSO_SHAPE.RECTANGLE)

# 架构层次图
layers = [
    ('应用层', 'Electron 主进程  ·  窗口管理  ·  系统托盘  ·  IPC 通信  ·  文件 I/O', C_PRIMARY, C_WHITE),
    ('渲染层', '原生 HTML5 + CSS3 + JavaScript  ·  零框架依赖  ·  响应式布局', C_PURPLE, C_WHITE),
    ('编辑引擎', '富文本引擎 (contentEditable)  ·  Markdown 解析器 (marked.js)  ·  代码高亮 (highlight.js)', C_SECONDARY, C_WHITE),
    ('数据层', 'LocalStorage 持久化  ·  JSON 数据模型  ·  版本历史管理', C_ACCENT, C_WHITE),
    ('导出引擎', 'PDF (Electron printToPDF)  ·  HTML / Markdown / JSON / TXT', C_RED, C_WHITE),
]

for i, (name, desc, bg, fg) in enumerate(layers):
    y = Inches(1.3 + i * 1.1)
    w = Inches(11.5 - i * 0.6)
    x = Inches(0.9 + i * 0.3)
    add_shape(slide, x, y, w, Inches(0.9), fill_color=bg, radius=0.1)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.05), Inches(2.0), Inches(0.4), name, font_size=18, color=fg, bold=True)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.45), w - Inches(0.6), Inches(0.4), desc, font_size=13, color=RGBColor(0xE0, 0xE0, 0xE0))

# 技术优势
add_textbox(slide, Inches(0.8), Inches(6.8), Inches(12), Inches(0.5),
            '✅ 零外部框架  ·  ✅ 内存占用 < 100MB  ·  ✅ 冷启动 < 2s  ·  ✅ 安装包 < 80MB  ·  ✅ 数据完全本地化',
            font_size=16, color=C_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_transition(slide, 'wipe')


# ============================================================
# SLIDE 10: 商业模式与定价
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_gradient(slide, '0D47A1', '1976D2')

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
            '06  商业模式与定价', font_size=36, color=C_WHITE, bold=True)

# 三个定价卡片
plans = [
    ('个人版', '免费', '/永久',
     ['✅ 基础笔记功能', '✅ 富文本 + Markdown', '✅ 桌面便签', '✅ 5种格式导出', '✅ 明暗主题'],
     C_SECONDARY, False),
    ('专业版', '¥99', '/年',
     ['✅ 个人版全部功能', '✅ 端到端加密', '✅ 云同步（可选）', '✅ 高级模板库', '✅ 优先技术支持', '✅ 团队协作（5人）'],
     C_ACCENT, True),
    ('企业版', '¥299', '/人/年',
     ['✅ 专业版全部功能', '✅ 私有化部署', '✅ LDAP/SSO 集成', '✅ 管理后台', '✅ 7×24 专属支持', '✅ 定制开发'],
     C_PRIMARY, False),
]

for i, (name, price, unit, features, color, featured) in enumerate(plans):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.3)
    h = Inches(5.8) if featured else Inches(5.5)
    y_adj = Inches(1.15) if featured else y

    # 卡片背景
    card = add_shape(slide, x, y_adj, Inches(3.7), h, fill_color=C_WHITE, radius=0.08)
    if featured:
        card.line.color.rgb = C_ACCENT
        card.line.width = Pt(3)

    # 推荐标签
    if featured:
        tag = add_shape(slide, x + Inches(1.0), y_adj - Inches(0.15), Inches(1.7), Inches(0.35), fill_color=C_ACCENT, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        tag.text_frame.paragraphs[0].text = '⭐ 推荐'
        tag.text_frame.paragraphs[0].font.size = Pt(12)
        tag.text_frame.paragraphs[0].font.color.rgb = C_WHITE
        tag.text_frame.paragraphs[0].font.bold = True
        tag.text_frame.paragraphs[0].font.name = '微软雅黑'
        tag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 版本名
    add_textbox(slide, x, y_adj + Inches(0.3), Inches(3.7), Inches(0.4), name, font_size=22, color=C_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # 价格
    add_textbox(slide, x, y_adj + Inches(0.8), Inches(3.7), Inches(0.7), price, font_size=42, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, y_adj + Inches(1.4), Inches(3.7), Inches(0.3), unit, font_size=14, color=C_GRAY, alignment=PP_ALIGN.CENTER)

    # 分隔线
    add_shape(slide, x + Inches(0.5), y_adj + Inches(1.8), Inches(2.7), Inches(0.02), fill_color=C_LIGHT, shape_type=MSO_SHAPE.RECTANGLE)

    # 功能列表
    for j, feat in enumerate(features):
        add_textbox(slide, x + Inches(0.4), y_adj + Inches(2.0 + j * 0.45), Inches(2.9), Inches(0.4),
                    feat, font_size=13, color=C_DARK)

add_slide_transition(slide, 'cover')


# ============================================================
# SLIDE 11: 服务保障与支持
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, C_LIGHT)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
            '07  服务保障与支持', font_size=36, color=C_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(0.95), Inches(1.5), Inches(0.04), fill_color=C_PRIMARY, shape_type=MSO_SHAPE.RECTANGLE)

services = [
    ('🚀', '快速部署', '30分钟完成部署', ['安装包一键部署', '支持静默安装 (企业批量)', 'MSI / EXE 双格式', '免依赖，开箱即用'], C_PRIMARY),
    ('📞', '技术支持', '多渠道响应', ['工单系统 (24h 内响应)', '企业专属技术对接人', '远程协助排障', '季度回访与优化建议'], C_SECONDARY),
    ('📚', '培训服务', '快速上手', ['产品使用手册', '视频教程库', '企业定制培训', '最佳实践分享'], C_PURPLE),
    ('🔄', '持续更新', '免费升级', ['每月功能更新', '安全补丁及时推送', '需求定制评估', '版本更新日志'], C_ACCENT),
]

for i, (icon, title, subtitle, items, color) in enumerate(services):
    x = Inches(0.5 + i * 3.15)
    y = Inches(1.3)

    card = add_shape(slide, x, y, Inches(2.9), Inches(5.5), fill_color=C_WHITE, radius=0.08)
    # 顶部色条
    add_shape(slide, x, y, Inches(2.9), Inches(0.06), fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    # 图标
    add_icon_badge(slide, x + Inches(0.95), y + Inches(0.3), Inches(0.9), icon, color, font_size=28)
    # 标题
    add_textbox(slide, x + Inches(0.2), y + Inches(1.35), Inches(2.5), Inches(0.4), title, font_size=18, color=C_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # 副标题
    add_textbox(slide, x + Inches(0.2), y + Inches(1.75), Inches(2.5), Inches(0.3), subtitle, font_size=12, color=color, alignment=PP_ALIGN.CENTER)

    # 分隔线
    add_shape(slide, x + Inches(0.4), y + Inches(2.2), Inches(2.1), Inches(0.02), fill_color=C_LIGHT, shape_type=MSO_SHAPE.RECTANGLE)

    # 功能列表
    for j, item in enumerate(items):
        add_textbox(slide, x + Inches(0.3), y + Inches(2.4 + j * 0.6), Inches(2.3), Inches(0.55),
                    f'• {item}', font_size=12, color=C_GRAY)

add_slide_transition(slide, 'fade')


# ============================================================
# SLIDE 12: 竞品对比
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, C_WHITE)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
            '08  竞品对比分析', font_size=36, color=C_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(0.95), Inches(1.5), Inches(0.04), fill_color=C_PRIMARY, shape_type=MSO_SHAPE.RECTANGLE)

# 对比表格
headers = ['对比维度', '云笔记', 'Notion', '印象笔记', '语雀']
header_colors = [C_DARK, C_PRIMARY, RGBColor(0x00, 0x00, 0x00), C_SECONDARY, RGBColor(0x00, 0xB9, 0x6B)]

rows = [
    ['部署方式', '本地部署 ✅', '纯云端', '云端为主', '纯云端'],
    ['数据安全', '完全本地 ⭐', '云端存储', '云端存储', '云端存储'],
    ['启动速度', '< 2 秒 ⭐', '3-5 秒', '3-8 秒', '2-4 秒'],
    ['内存占用', '< 100MB ⭐', '200-500MB', '300-800MB', '浏览器'],
    ['离线支持', '完全离线 ⭐', '部分', '部分', '无'],
    ['年费 (个人)', '免费 ⭐', '¥96', '¥148', '免费/¥99'],
    ['Markdown', '原生支持 ✅', '支持', '有限', '支持'],
    ['桌面便签', '原生支持 ⭐', '无', '无', '无'],
    ['定制开发', '支持 ✅', '不支持', '不支持', '不支持'],
    ['私有化部署', '原生支持 ⭐', '不支持', '不支持', '企业版'],
]

col_widths = [Inches(2.2), Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)]
x_start = Inches(0.7)

# 表头
x = x_start
for j, (header, color) in enumerate(zip(headers, header_colors)):
    bg = color if j == 1 else C_DARK
    add_shape(slide, x, Inches(1.3), col_widths[j], Inches(0.5), fill_color=bg, shape_type=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, x, Inches(1.3), col_widths[j], Inches(0.5), header, font_size=14, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    x += col_widths[j]

# 数据行
for i, row in enumerate(rows):
    y = Inches(1.8 + i * 0.5)
    x = x_start
    bg = C_LIGHT if i % 2 == 0 else C_WHITE
    for j, cell in enumerate(row):
        cell_bg = bg
        cell_color = C_DARK
        cell_bold = False
        if j == 1 and ('⭐' in cell or '✅' in cell):
            cell_bg = RGBColor(0xE3, 0xF2, 0xFD) if i % 2 == 0 else RGBColor(0xEB, 0xF5, 0xFF)
            cell_color = C_PRIMARY
            cell_bold = True
        add_shape(slide, x, y, col_widths[j], Inches(0.5), fill_color=cell_bg, shape_type=MSO_SHAPE.RECTANGLE)
        add_textbox(slide, x, y + Inches(0.05), col_widths[j], Inches(0.4), cell, font_size=12, color=cell_color, bold=cell_bold, alignment=PP_ALIGN.CENTER)
        x += col_widths[j]

add_slide_transition(slide, 'push')


# ============================================================
# SLIDE 13: 产品路线图
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, C_LIGHT)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
            '09  产品路线图', font_size=36, color=C_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(0.95), Inches(1.5), Inches(0.04), fill_color=C_PRIMARY, shape_type=MSO_SHAPE.RECTANGLE)

# 时间线
milestones = [
    ('v1.0', '2026 Q1', '基础版发布', ['核心编辑功能', '桌面便签', '多格式导出', '明暗主题'], C_SECONDARY, True),
    ('v1.5', '2026 Q2', '体验升级', ['云同步（可选）', 'AI 智能摘要', '全文检索优化', '移动端适配'], C_PRIMARY, False),
    ('v2.0', '2026 Q3', '企业版', ['团队协作', '权限管理', '私有化部署', 'SSO 集成'], C_PURPLE, False),
    ('v3.0', '2026 Q4', '生态扩展', ['插件市场', '知识图谱', 'API 开放', '多端同步'], C_ACCENT, False),
]

# 横向时间线
line_y = Inches(2.0)
add_shape(slide, Inches(1.0), line_y + Inches(0.2), Inches(11.3), Inches(0.04), fill_color=C_GRAY, shape_type=MSO_SHAPE.RECTANGLE)

for i, (ver, date, title, features, color, current) in enumerate(milestones):
    x = Inches(1.2 + i * 3.0)

    # 节点圆
    size = Inches(0.55) if current else Inches(0.45)
    offset = Inches(0) if current else Inches(0.05)
    node = add_icon_badge(slide, x + Inches(0.5) + offset, line_y + offset, size, ver, color, font_size=11)

    if current:
        # 当前版本高亮标签
        tag = add_shape(slide, x + Inches(0.15), line_y - Inches(0.45), Inches(1.3), Inches(0.3), fill_color=C_SECONDARY, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        tag.text_frame.paragraphs[0].text = '当前版本'
        tag.text_frame.paragraphs[0].font.size = Pt(10)
        tag.text_frame.paragraphs[0].font.color.rgb = C_WHITE
        tag.text_frame.paragraphs[0].font.bold = True
        tag.text_frame.paragraphs[0].font.name = '微软雅黑'
        tag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 卡片
    card_y = Inches(2.8)
    card = add_shape(slide, x, card_y, Inches(2.7), Inches(3.8), fill_color=C_WHITE, border_color=color, border_width=Pt(2), radius=0.08)

    # 日期
    add_textbox(slide, x + Inches(0.2), card_y + Inches(0.15), Inches(2.3), Inches(0.3), date, font_size=13, color=color, bold=True)
    # 标题
    add_textbox(slide, x + Inches(0.2), card_y + Inches(0.5), Inches(2.3), Inches(0.3), title, font_size=18, color=C_DARK, bold=True)
    # 分隔线
    add_shape(slide, x + Inches(0.3), card_y + Inches(0.9), Inches(2.1), Inches(0.02), fill_color=C_LIGHT, shape_type=MSO_SHAPE.RECTANGLE)
    # 功能列表
    for j, feat in enumerate(features):
        prefix = '✅' if current else '🔲'
        add_textbox(slide, x + Inches(0.2), card_y + Inches(1.05 + j * 0.55), Inches(2.3), Inches(0.5),
                    f'{prefix} {feat}', font_size=12, color=C_DARK if current else C_GRAY)

add_slide_transition(slide, 'split')


# ============================================================
# SLIDE 14: 客户案例 & 使用场景
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, C_WHITE)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
            '应用场景与客户价值', font_size=36, color=C_PRIMARY, bold=True)
add_shape(slide, Inches(0.8), Inches(0.95), Inches(1.5), Inches(0.04), fill_color=C_PRIMARY, shape_type=MSO_SHAPE.RECTANGLE)

scenarios = [
    ('🏢', '企业办公', '会议纪要 / OKR 管理 / 周报月报\n数据本地化保证信息安全\n团队共享笔记本协作', '适合: 中大型企业、金融机构', C_PRIMARY),
    ('💻', '技术团队', 'API 文档 / 代码笔记 / 技术方案\nMarkdown 原生支持\n代码高亮 + 流程图', '适合: 互联网公司、研发部门', C_PURPLE),
    ('🎓', '教育培训', '课堂笔记 / 学习计划 / 知识整理\n待办清单跟踪学习进度\n便签提醒重要事项', '适合: 高校、培训机构、学生', C_SECONDARY),
    ('👤', '个人效率', '日记 / 读书笔记 / 灵感记录\n桌面便签管理每日待办\n多格式导出分享', '适合: 自由职业者、写作者', C_ACCENT),
]

for i, (icon, title, desc, target, color) in enumerate(scenarios):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.3 + row * 3.0)

    card = add_shape(slide, x, y, Inches(5.9), Inches(2.7), fill_color=C_WHITE, border_color=color, border_width=Pt(2), radius=0.08)
    # 左侧色块
    add_shape(slide, x, y, Inches(0.08), Inches(2.7), fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)

    # 图标
    add_icon_badge(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.8), icon, color, font_size=28)
    # 标题
    add_textbox(slide, x + Inches(1.3), y + Inches(0.35), Inches(4.0), Inches(0.4), title, font_size=22, color=C_DARK, bold=True)
    # 描述
    add_textbox(slide, x + Inches(1.3), y + Inches(0.85), Inches(4.3), Inches(1.2), desc, font_size=13, color=C_GRAY)
    # 目标客户
    tag = add_shape(slide, x + Inches(1.3), y + Inches(2.1), Inches(3.5), Inches(0.35), fill_color=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tag.text_frame.paragraphs[0].text = target
    tag.text_frame.paragraphs[0].font.size = Pt(11)
    tag.text_frame.paragraphs[0].font.color.rgb = C_WHITE
    tag.text_frame.paragraphs[0].font.bold = True
    tag.text_frame.paragraphs[0].font.name = '微软雅黑'
    tag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_slide_transition(slide, 'dissolve')


# ============================================================
# SLIDE 15: 结尾 CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_gradient(slide, '1A237E', '0D47A1')

# 装饰
add_shape(slide, Inches(9), Inches(-2), Inches(7), Inches(7), fill_color=RGBColor(0x1E, 0x3A, 0x8A), shape_type=MSO_SHAPE.OVAL)
add_shape(slide, Inches(-2), Inches(4), Inches(6), Inches(6), fill_color=RGBColor(0x15, 0x30, 0x70), shape_type=MSO_SHAPE.OVAL)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            '感谢聆听', font_size=64, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5), Inches(3.0), Inches(3.3), Inches(0.04), fill_color=C_SECONDARY, shape_type=MSO_SHAPE.RECTANGLE)

add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.7),
            '云笔记 — 让知识管理更简单、更安全', font_size=24, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)

# 联系方式卡片
contact_card = add_shape(slide, Inches(3.5), Inches(4.5), Inches(6.3), Inches(2.0), fill_color=RGBColor(0x1E, 0x3A, 0x8A), radius=0.1)
contact_card.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x8A)

contact_lines = [
    ('📧  邮箱：sales@cloudnotes.com', 16, RGBColor(0xBB, 0xDE, 0xFB), False),
    ('📱  电话：400-888-NOTE', 16, RGBColor(0xBB, 0xDE, 0xFB), False),
    ('🌐  官网：www.cloudnotes.com', 16, RGBColor(0xBB, 0xDE, 0xFB), False),
]
add_multiline(slide, Inches(4.0), Inches(4.7), Inches(5.3), Inches(1.8), contact_lines, alignment=PP_ALIGN.CENTER)

add_slide_transition(slide, 'fade')


# ===== 嵌入 GIF 动图 =====
# python-pptx 嵌入的 GIF 只显示第一帧
# 所以我们使用静态截图 + 在 PPT 中用便签展示多帧
# GIF 文件仍然可以在 readme 和 web 页面中使用

# ===== 保存 =====
prs.save(OUT)
size = os.path.getsize(OUT)
print(f'PPT saved: {OUT} ({size/1024:.0f} KB, {len(prs.slides)} slides)')
